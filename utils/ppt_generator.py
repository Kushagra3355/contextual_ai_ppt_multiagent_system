from pptx import Presentation
from typing import List
from orchestrator.agent_state import SlideValidation


def create_presentation(
    validation_results: List[SlideValidation], output_path: str, title: str
) -> str:
    """Create PowerPoint presentation from validation results"""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "AI Generated Presentation"

    # Content slides
    for slide_validation in validation_results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_validation.title

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for i, validation in enumerate(slide_validation.validation):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"• {validation.point}"

            # Add validation status in parentheses with italic
            run = p.add_run()
            run.text = f" ({validation.status})"
            run.font.italic = True

    prs.save(output_path)
    return output_path
